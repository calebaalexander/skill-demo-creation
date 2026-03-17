#!/usr/bin/env python3
"""Generate a branded Snowflake PPTX deck by cloning template slides.

Usage:
    uv run generate_pptx.py \
        --template path/to/SNOWFLAKE_TEMPLATE_JANUARY_2026.pptx \
        --slides-json slides_content.json \
        --output output_deck.pptx

The slides-json file contains an array of slide definitions. Each object has:
    {
        "type": "cover" | "safe_harbor" | "agenda" | "chapter" | "content_1col"
                | "content_2col" | "content_2col_titled" | "content_3col_titled"
                | "content_3col_icons" | "content_4col_icons"
                | "speakers" | "template_clone" | "thank_you",
        "title_line1": "BLACK LINE",          # cover/chapter/thank_you colorstack
        "title_line2": "WHITE LINE",          # cover/chapter/thank_you colorstack
        "title": "Slide Title",               # content slides
        "subtitle": "Optional subtitle",
        "body": "Body text",                  # 1-col content
        "left_title": "Left Header",          # 2-col titled
        "left_body": "Left text",
        "right_title": "Right Header",
        "right_body": "Right text",
        "col1_title": "", "col1_body": "",    # 3-col titled / icons
        "col2_title": "", "col2_body": "",
        "col3_title": "", "col3_body": "",
        "col4_title": "", "col4_body": "",    # 4-col icons
        "col1_icon": "", "col2_icon": "",     # per-column icon labels
        "col3_icon": "", "col4_icon": "",
        "template_index": 0,                    # template_clone
        "agenda_items": ["Item 1", "Item 2"], # agenda
        "presenter": "Name",                  # cover
        "date": "March 2026",                 # cover
        "speakers": [                          # speakers slide
            {"name": "First Last", "title": "Title", "email": "first.last@snowflake.com",
             "photo": "optional/path/or/url"}
        ],
        "speaker_notes": "Notes text"
    }
"""

import argparse
import copy
import json
import sys
from pathlib import Path

import os
import tempfile
import urllib.request
import urllib.parse

from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

SNOWFLAKE_BLUE = RGBColor(0x29, 0xB5, 0xE8)
MID_BLUE = RGBColor(0x11, 0x56, 0x7F)
MEDIUM_GRAY = RGBColor(0x5B, 0x5B, 0x5B)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

TEMPLATE_SLIDE_MAP = {
    "cover": 19,
    "cover_logo": 21,
    "safe_harbor": 2,
    "chapter_1": 22,
    "chapter_2": 23,
    "chapter_3": 24,
    "content_1col": 32,
    "content_2col": 33,
    "content_2col_titled": 34,
    "content_3col_titled": 36,
    "content_3col_icons": 37,
    "content_4col_icons": 39,
    "agenda": 44,
    "speakers": 32,
    "thank_you": 72,
}


def clone_slide(prs, source_slide_index):
    template_slides = list(prs.slides)
    source_slide = template_slides[source_slide_index]
    slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    for shape in source_slide.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)

    return new_slide


def set_run_text(run_element, text, color=None, size=None, bold=None):
    run_element.text = text
    rPr = run_element.find(qn("a:rPr"))
    if rPr is None:
        rPr = run_element.makeelement(qn("a:rPr"), {})
        run_element.insert(0, rPr)

    if color is not None:
        solidFill = rPr.find(qn("a:solidFill"))
        if solidFill is not None:
            rPr.remove(solidFill)
        solidFill = rPr.makeelement(qn("a:solidFill"), {})
        srgbClr = solidFill.makeelement(
            qn("a:srgbClr"), {"val": "%02X%02X%02X" % (color[0], color[1], color[2])}
        )
        solidFill.append(srgbClr)
        rPr.append(solidFill)

        for scheme_fill in rPr.findall(qn("a:schemeClr")):
            rPr.remove(scheme_fill)

    if size is not None:
        rPr.set("sz", str(int(size.pt * 100)))

    if bold is not None:
        rPr.set("b", "1" if bold else "0")


def find_shapes_with_text(slide):
    result = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            result.append(shape)
    return result


def get_all_runs(shape):
    runs = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            runs.append(run)
    return runs


def set_shape_text(shape, text, color=None, size=None, bold=None):
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.text = ""
    first_run = None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            first_run = run
            break
        if first_run:
            break
    if first_run:
        first_run.text = text
        if color:
            set_run_text(first_run._r, text, color, size, bold)


def set_body_text_preserve_spacer(shape, text):
    paras = shape.text_frame.paragraphs
    if len(paras) >= 2 and paras[0].text.strip() == "":
        for para in paras[1:]:
            for run in para.runs:
                run.text = ""
        if paras[1].runs:
            paras[1].runs[0].text = text
    else:
        set_shape_text(shape, text)


def remove_shape(slide, shape):
    sp = shape._element
    sp.getparent().remove(sp)


def constrain_body_autofit(shape):
    if not shape.has_text_frame:
        return
    bodyPr = shape.text_frame._txBody.find(qn("a:bodyPr"))
    if bodyPr is None:
        return
    for child_tag in ("a:spAutoFit", "a:noAutofit", "a:normAutofit"):
        el = bodyPr.find(qn(child_tag))
        if el is not None:
            bodyPr.remove(el)
    bodyPr.append(bodyPr.makeelement(qn("a:normAutofit"), {"fontScale": "90000", "lnSpcReduction": "10000"}))


def constrain_title_autofit(shape):
    if not shape.has_text_frame:
        return
    bodyPr = shape.text_frame._txBody.find(qn("a:bodyPr"))
    if bodyPr is None:
        return
    for child_tag in ("a:spAutoFit", "a:noAutofit"):
        el = bodyPr.find(qn(child_tag))
        if el is not None:
            bodyPr.remove(el)
    if bodyPr.find(qn("a:normAutofit")) is None:
        bodyPr.append(bodyPr.makeelement(qn("a:normAutofit"), {"fontScale": "80000", "lnSpcReduction": "10000"}))


def add_body_top_inset(shape, inset_emu=500000):
    if not shape.has_text_frame:
        return
    bodyPr = shape.text_frame._txBody.find(qn("a:bodyPr"))
    if bodyPr is None:
        return
    bodyPr.set("tIns", str(inset_emu))


def expand_shape_height(shape, target_bottom_emu):
    new_h = int(target_bottom_emu) - shape.top
    if new_h > shape.height:
        shape.height = new_h


def _split_body_sentences(text):
    import re
    parts = re.split(r'(?<=\.)\s+', text.strip())
    return [p.strip() for p in parts if p.strip()]


def set_body_bullets(shape, text, font_size_hundredths=1400, line_spacing_pct=115000, top_inset=None):
    if not shape.has_text_frame:
        return
    if top_inset is not None:
        bodyPr = shape.text_frame._txBody.find(qn("a:bodyPr"))
        if bodyPr is not None:
            bodyPr.set("tIns", str(int(top_inset)))
    txBody = shape.text_frame._txBody
    existing = txBody.findall(qn("a:p"))
    if not existing:
        return
    ref_p = existing[0]
    ref_rPr = None
    for run in ref_p.findall(qn("a:r")):
        rp = run.find(qn("a:rPr"))
        if rp is not None:
            ref_rPr = copy.deepcopy(rp)
            break
    for p in existing:
        txBody.remove(p)
    sentences = _split_body_sentences(text)
    for i, sentence in enumerate(sentences):
        p_el = txBody.makeelement(qn("a:p"), {})
        pPr = p_el.makeelement(qn("a:pPr"), {})
        buChar = pPr.makeelement(qn("a:buChar"), {"char": "\u2022"})
        pPr.append(buChar)
        lnSpc = pPr.makeelement(qn("a:lnSpc"), {})
        spcPct = lnSpc.makeelement(qn("a:spcPct"), {"val": str(line_spacing_pct)})
        lnSpc.append(spcPct)
        pPr.append(lnSpc)
        if i > 0:
            spcBef = pPr.makeelement(qn("a:spcBef"), {})
            spcPts = spcBef.makeelement(qn("a:spcPts"), {"val": "600"})
            spcBef.append(spcPts)
            pPr.append(spcBef)
        p_el.append(pPr)
        r_el = p_el.makeelement(qn("a:r"), {})
        if ref_rPr is not None:
            new_rPr = copy.deepcopy(ref_rPr)
            new_rPr.set("sz", str(font_size_hundredths))
        else:
            new_rPr = r_el.makeelement(qn("a:rPr"), {"lang": "en-US", "sz": str(font_size_hundredths)})
        r_el.append(new_rPr)
        t_el = r_el.makeelement(qn("a:t"), {})
        t_el.text = sentence
        r_el.append(t_el)
        p_el.append(r_el)
        txBody.append(p_el)


ICON_SIZE = Emu(1500000)
ICON_MARGIN_RIGHT = Emu(457200)
ICON_MARGIN_BOTTOM = Emu(457200)


def _build_icon_index(template_prs):
    index = {}
    template_slides = list(template_prs.slides)
    for slide_idx in range(63, min(72, len(template_slides))):
        slide = template_slides[slide_idx]
        text_shapes = []
        group_shapes = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.top > Emu(900000):
                t = shape.text_frame.text.strip().replace("\n", " ").replace("  ", " ")
                if t and t != "\u2039#\u203a":
                    text_shapes.append((shape.top, shape.left, shape.width, shape.height, t))
            if shape.shape_type == 6:
                group_shapes.append(shape)
        for ts_top, ts_left, ts_w, ts_h, ts_text in text_shapes:
            ts_center_x = ts_left + ts_w // 2
            best_group = None
            best_dist = float("inf")
            for gs in group_shapes:
                gs_center_x = gs.left + gs.width // 2
                dx = abs(gs_center_x - ts_center_x)
                if gs.top < ts_top and dx < Emu(500000):
                    dist = abs(gs.top + gs.height - ts_top) + dx
                    if dist < best_dist:
                        best_dist = dist
                        best_group = gs
            if best_group:
                key = ts_text.lower()
                index[key] = best_group
    return index


ICON_LABEL_OVERRIDES = {
    "self-service": "users",
}


def _find_icon(icon_index, label):
    label_lower = label.lower().strip()
    resolved = ICON_LABEL_OVERRIDES.get(label_lower, label_lower)
    if resolved in icon_index:
        return icon_index[resolved]
    for key, group in icon_index.items():
        if resolved in key:
            return group
    if resolved != label_lower:
        if label_lower in icon_index:
            return icon_index[label_lower]
    return None


def place_icon(slide, icon_group, left, top, size=ICON_SIZE):
    el = copy.deepcopy(icon_group._element)
    orig_w = icon_group.width
    orig_h = icon_group.height
    scale = min(size / orig_w, size / orig_h) if orig_w and orig_h else 1.0
    new_w = int(orig_w * scale)
    new_h = int(orig_h * scale)
    xfrm = el.find(".//" + qn("a:xfrm"))
    if xfrm is None:
        grpSpPr = el.find(qn("p:grpSpPr"))
        if grpSpPr is not None:
            xfrm = grpSpPr.find(qn("a:xfrm"))
    if xfrm is not None:
        off = xfrm.find(qn("a:off"))
        ext = xfrm.find(qn("a:ext"))
        if off is not None:
            off.set("x", str(int(left)))
            off.set("y", str(int(top)))
        if ext is not None:
            ext.set("cx", str(new_w))
            ext.set("cy", str(new_h))
    slide.shapes._spTree.append(el)


def set_colorstack(shape, line1, line2):
    runs = get_all_runs(shape)
    if len(runs) >= 2:
        runs[0].text = line1
        runs[1].text = line2
        for r in runs[2:]:
            r.text = ""


def set_speaker_notes(slide, notes_text):
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


SLIDE_WIDTH = Emu(9144000)
RIGHT_MARGIN = Emu(366600)
TOP_MARGIN = Emu(300000)
LOGO_MAX_HEIGHT = Emu(365000)
LOGO_MAX_WIDTH = Emu(1600000)
SF_LOGO_VCENTER = Emu(757646)


def add_logo_to_slide(slide, logo_path, left, top, width=None, height=None):
    if not logo_path or not os.path.exists(logo_path):
        return
    if width:
        slide.shapes.add_picture(logo_path, left, top, width=width)
    elif height:
        slide.shapes.add_picture(logo_path, left, top, height=height)
    else:
        slide.shapes.add_picture(logo_path, left, top)


def place_logo_top_right(slide, logo_path):
    if not logo_path or not os.path.exists(logo_path):
        return
    from PIL import Image as _Image
    _img = _Image.open(logo_path)
    _w, _h = _img.size
    aspect = _w / _h
    logo_height = LOGO_MAX_HEIGHT
    logo_width = int(logo_height * aspect)
    if logo_width > LOGO_MAX_WIDTH:
        logo_width = LOGO_MAX_WIDTH
        logo_height = int(logo_width / aspect)
    logo_left = SLIDE_WIDTH - RIGHT_MARGIN - logo_width
    logo_top = int(SF_LOGO_VCENTER) - logo_height // 2
    slide.shapes.add_picture(logo_path, logo_left, logo_top, width=logo_width, height=logo_height)


def apply_cover(slide, data):
    shapes = find_shapes_with_text(slide)
    shapes_by_top = sorted(shapes, key=lambda s: s.top)
    for shape in shapes_by_top:
        runs = get_all_runs(shape)
        if not runs:
            continue
        if shape.top < Emu(2000000) and len(runs) >= 2:
            set_colorstack(shape, data.get("title_line1", ""), data.get("title_line2", ""))
        elif Emu(2800000) < shape.top < Emu(3800000):
            set_shape_text(shape, data.get("subtitle", ""))
        elif shape.top > Emu(4000000) and shape.height < Emu(500000):
            presenter = data.get("presenter", "")
            date = data.get("date", "")
            set_shape_text(shape, f"{presenter}  |  {date}")

    set_speaker_notes(slide, data.get("speaker_notes", ""))


def apply_chapter(slide, data):
    shapes = find_shapes_with_text(slide)
    for shape in shapes:
        runs = get_all_runs(shape)
        if len(runs) >= 2 and shape.width > Emu(5000000):
            set_colorstack(shape, data.get("title_line1", ""), data.get("title_line2", ""))
            break

    set_speaker_notes(slide, data.get("speaker_notes", ""))


FOOTER_TOP = Emu(4600000)


BODY_FULL_WIDTH = Emu(8344800)
BODY_WITH_ICON_WIDTH = Emu(6261300)


def apply_content_1col(slide, data):
    shapes = find_shapes_with_text(slide)
    title_set = False
    subtitle_set = False
    body_set = False
    body_shape = None
    shapes_to_remove = []
    has_icon = bool(data.get("icon"))

    shapes_by_top = sorted(shapes, key=lambda s: s.top)

    for shape in shapes_by_top:
        page_num = shape.text_frame.text.strip()
        if page_num == '\u2039#\u203a' or (shape.width < Emu(600000) and shape.top > Emu(4500000)):
            continue

        if shape.top < Emu(500000) and shape.width > Emu(4000000) and not title_set:
            set_shape_text(shape, data.get("title", ""))
            constrain_title_autofit(shape)
            title_set = True
        elif Emu(500000) <= shape.top < Emu(1200000) and not subtitle_set:
            subtitle_text = data.get("subtitle", "")
            if subtitle_text:
                set_shape_text(shape, subtitle_text)
                constrain_title_autofit(shape)
            else:
                shapes_to_remove.append(shape)
            subtitle_set = True
        elif shape.top >= Emu(1200000) and shape.width > Emu(4000000) and not body_set:
            body_shape = shape
            if has_icon:
                shape.width = int(BODY_WITH_ICON_WIDTH)
            else:
                shape.width = int(BODY_FULL_WIDTH)
            expand_shape_height(shape, FOOTER_TOP)
            set_body_bullets(shape, data.get("body", ""))
            constrain_body_autofit(shape)
            body_set = True

    for s in shapes_to_remove:
        remove_shape(slide, s)

    set_speaker_notes(slide, data.get("speaker_notes", ""))


def apply_content_2col_titled(slide, data):
    shapes = find_shapes_with_text(slide)

    title_shape = None
    subtitle_shape = None
    para_titles = []
    body_shapes = []

    for shape in shapes:
        page_num = shape.text_frame.text.strip()
        if page_num == '\u2039#\u203a' or (shape.width < Emu(600000) and shape.top > Emu(4500000)):
            continue

        if shape.top < Emu(500000) and shape.width > Emu(4000000):
            title_shape = shape
        elif Emu(500000) <= shape.top < Emu(1200000) and shape.width > Emu(4000000):
            subtitle_shape = shape
        elif shape.top >= Emu(1200000) and shape.width < Emu(4000000) and shape.height < Emu(600000):
            para_titles.append(shape)
        elif shape.top >= Emu(1200000) and shape.height > Emu(600000):
            body_shapes.append(shape)

    if title_shape:
        set_shape_text(title_shape, data.get("title", ""))
        constrain_title_autofit(title_shape)
    if subtitle_shape:
        subtitle_text = data.get("subtitle", "")
        if subtitle_text:
            set_shape_text(subtitle_shape, subtitle_text)
            constrain_title_autofit(subtitle_shape)
        else:
            remove_shape(slide, subtitle_shape)

    para_titles.sort(key=lambda s: s.left)
    body_shapes.sort(key=lambda s: s.left)

    if len(para_titles) >= 1:
        set_shape_text(para_titles[0], data.get("left_title", ""))
        constrain_title_autofit(para_titles[0])
    if len(para_titles) >= 2:
        set_shape_text(para_titles[1], data.get("right_title", ""))
        constrain_title_autofit(para_titles[1])
    col_title_inset = Emu(500000)
    if len(body_shapes) >= 1:
        set_body_bullets(body_shapes[0], data.get("left_body", ""), font_size_hundredths=1100, top_inset=col_title_inset)
        constrain_body_autofit(body_shapes[0])
    if len(body_shapes) >= 2:
        set_body_bullets(body_shapes[1], data.get("right_body", ""), font_size_hundredths=1100, top_inset=col_title_inset)
        constrain_body_autofit(body_shapes[1])

    set_speaker_notes(slide, data.get("speaker_notes", ""))


def apply_content_3col_titled(slide, data):
    shapes = find_shapes_with_text(slide)

    title_shape = None
    subtitle_shape = None
    col_titles = []
    col_bodies = []

    for shape in shapes:
        page_num = shape.text_frame.text.strip()
        if page_num == '\u2039#\u203a' or (shape.width < Emu(600000) and shape.top > Emu(4500000)):
            continue

        if shape.top < Emu(400000) and shape.width > Emu(4000000):
            title_shape = shape
        elif Emu(400000) <= shape.top < Emu(1200000) and shape.width > Emu(4000000):
            subtitle_shape = shape
        elif shape.top >= Emu(1200000) and shape.height < Emu(600000):
            col_titles.append(shape)
        elif shape.top >= Emu(1200000) and shape.height >= Emu(600000):
            col_bodies.append(shape)

    if title_shape:
        set_shape_text(title_shape, data.get("title", ""))
        constrain_title_autofit(title_shape)
    if subtitle_shape:
        subtitle_text = data.get("subtitle", "")
        if subtitle_text:
            set_shape_text(subtitle_shape, subtitle_text)
            constrain_title_autofit(subtitle_shape)
        else:
            remove_shape(slide, subtitle_shape)

    col_titles.sort(key=lambda s: s.left)
    col_bodies.sort(key=lambda s: s.left)

    for i, prefix in enumerate(["col1", "col2", "col3"]):
        if i < len(col_titles):
            set_shape_text(col_titles[i], data.get(f"{prefix}_title", ""))
            constrain_title_autofit(col_titles[i])
        if i < len(col_bodies):
            set_body_bullets(col_bodies[i], data.get(f"{prefix}_body", ""), font_size_hundredths=1100, top_inset=Emu(500000))
            constrain_body_autofit(col_bodies[i])

    set_speaker_notes(slide, data.get("speaker_notes", ""))


ICON_COL_SIZE = Emu(700000)


def _apply_col_icons(slide, data, col_count, icon_index):
    shapes = list(slide.shapes)

    title_shape = None
    subtitle_shape = None
    freeform_icons = []
    para_titles = []
    body_shapes = []

    for shape in shapes:
        if shape.shape_type == 5 and shape.top >= Emu(1300000) and shape.top < Emu(1700000):
            freeform_icons.append(shape)
        elif shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt == '\u2039#\u203a' or (shape.width < Emu(600000) and shape.top > Emu(4500000)):
                continue
            if shape.top < Emu(500000) and shape.width > Emu(4000000):
                title_shape = shape
            elif Emu(500000) <= shape.top < Emu(1200000) and shape.width > Emu(4000000):
                subtitle_shape = shape
            elif shape.top >= Emu(1800000) and shape.height < Emu(600000) and shape.width < Emu(3000000):
                para_titles.append(shape)
            elif shape.top >= Emu(1200000) and shape.height > Emu(600000):
                body_shapes.append(shape)

    if title_shape:
        set_shape_text(title_shape, data.get("title", ""))
        constrain_title_autofit(title_shape)
    if subtitle_shape:
        subtitle_text = data.get("subtitle", "")
        if subtitle_text:
            set_shape_text(subtitle_shape, subtitle_text)
            constrain_title_autofit(subtitle_shape)
        else:
            remove_shape(slide, subtitle_shape)

    freeform_icons.sort(key=lambda s: s.left)
    for fi in freeform_icons:
        remove_shape(slide, fi)

    para_titles.sort(key=lambda s: s.left)
    body_shapes.sort(key=lambda s: s.left)

    prefixes = [f"col{i+1}" for i in range(col_count)]
    for i, prefix in enumerate(prefixes):
        icon_label = data.get(f"{prefix}_icon", "")
        if icon_label and icon_index:
            icon_group = _find_icon(icon_index, icon_label)
            if icon_group and i < len(para_titles):
                pt = para_titles[i]
                icon_cx = int(ICON_COL_SIZE)
                col_center_x = pt.left + pt.width // 2
                icon_left = col_center_x - icon_cx // 2
                icon_top = Emu(1400000)
                place_icon(slide, icon_group, icon_left, icon_top, size=ICON_COL_SIZE)

        if i < len(para_titles):
            set_shape_text(para_titles[i], data.get(f"{prefix}_title", ""))
            constrain_title_autofit(para_titles[i])
        if i < len(body_shapes):
            body_text = data.get(f"{prefix}_body", "")
            set_body_bullets(body_shapes[i], body_text, font_size_hundredths=1100, top_inset=Emu(500000))
            constrain_body_autofit(body_shapes[i])

    set_speaker_notes(slide, data.get("speaker_notes", ""))


def apply_content_3col_icons(slide, data, icon_index=None):
    _apply_col_icons(slide, data, 3, icon_index)


def apply_content_4col_icons(slide, data, icon_index=None):
    _apply_col_icons(slide, data, 4, icon_index)


def apply_template_clone(slide, data):
    set_speaker_notes(slide, data.get("speaker_notes", ""))


def apply_agenda(slide, data):
    shapes = find_shapes_with_text(slide)
    items = data.get("agenda_items", [])

    text_shape = None
    for shape in shapes:
        if shape.width > Emu(4000000) and shape.top > Emu(800000):
            text_shape = shape
            break

    if text_shape:
        bodyPr = text_shape.text_frame._txBody.find(qn("a:bodyPr"))

        paras = list(text_shape.text_frame.paragraphs)
        for i, para in enumerate(paras):
            if i < len(items):
                for run in para.runs:
                    run.text = items[i]
                pPr = para._p.find(qn("a:pPr"))
                if pPr is None:
                    pPr = para._p.makeelement(qn("a:pPr"), {})
                    para._p.insert(0, pPr)
                lnSpc = pPr.find(qn("a:lnSpc"))
                if lnSpc is None:
                    lnSpc = pPr.makeelement(qn("a:lnSpc"), {})
                    pPr.append(lnSpc)
                for child in list(lnSpc):
                    lnSpc.remove(child)
                spcPct = lnSpc.makeelement(qn("a:spcPct"), {"val": "115000"})
                lnSpc.append(spcPct)
            else:
                p_elem = para._p
                p_elem.getparent().remove(p_elem)

    arrows = []
    for shape in slide.shapes:
        if shape.width < Emu(200000) and shape.top > Emu(1000000) and shape.top < Emu(4800000):
            arrows.append(shape)
    arrows.sort(key=lambda s: s.top)

    for i in range(len(arrows) - 1, -1, -1):
        if i >= len(items):
            sp = arrows[i]._element
            sp.getparent().remove(sp)
    arrows = [a for a in arrows if a._element.getparent() is not None]

    if text_shape and arrows and items:
        ARROW_SPACING = 443089
        shape_top = text_shape.top
        shape_h = text_shape.height
        n = len(items)
        total_block = ARROW_SPACING * (n - 1)
        shape_center_y = shape_top + shape_h // 2
        first_arrow_cy = shape_center_y - total_block // 2
        for i, arrow in enumerate(arrows):
            cy = first_arrow_cy + i * ARROW_SPACING
            arrow.top = int(cy - arrow.height // 2)

    set_speaker_notes(slide, data.get("speaker_notes", ""))


def apply_thank_you(slide, data):
    shapes = find_shapes_with_text(slide)
    for shape in shapes:
        runs = get_all_runs(shape)
        if any("THANK" in r.text.upper() or "YOU" in r.text.upper() for r in runs):
            set_colorstack(shape, data.get("title_line1", "THANK"), data.get("title_line2", "YOU"))
        elif shape.top > Emu(2500000):
            contact = data.get("contact", "")
            set_shape_text(shape, contact)

    set_speaker_notes(slide, data.get("speaker_notes", ""))


def apply_cover_logo(slide, data):
    shapes = find_shapes_with_text(slide)
    shapes_by_top = sorted(shapes, key=lambda s: s.top)
    for shape in shapes_by_top:
        runs = get_all_runs(shape)
        if not runs:
            continue
        if shape.top < Emu(2500000) and shape.width > Emu(4000000) and len(runs) >= 2:
            set_colorstack(shape, data.get("title_line1", ""), data.get("title_line2", ""))
        elif Emu(2800000) < shape.top < Emu(4200000) and shape.width > Emu(4000000):
            set_shape_text(shape, data.get("subtitle", ""))
        elif shape.top > Emu(4200000) and shape.height < Emu(500000) and shape.width > Emu(2000000):
            presenter = data.get("presenter", "")
            date = data.get("date", "")
            set_shape_text(shape, f"{presenter}  |  {date}")

    set_speaker_notes(slide, data.get("speaker_notes", ""))


def _generate_initials_avatar(name, dest, size=400):
    from PIL import Image as _Img, ImageDraw as _Draw, ImageFont as _Font
    img = _Img.new("RGBA", (size, size), (0, 0, 0, 0))
    draw = _Draw.Draw(img)
    draw.ellipse((0, 0, size, size), fill=(41, 181, 232, 255))
    initials = "".join(p[0].upper() for p in name.strip().split() if p)[:2]
    try:
        font = _Font.truetype("/System/Library/Fonts/Helvetica.ttc", size // 3)
    except Exception:
        try:
            font = _Font.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", size // 3)
        except Exception:
            font = _Font.load_default()
    bbox = draw.textbbox((0, 0), initials, font=font)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    draw.text(((size - tw) / 2, (size - th) / 2 - bbox[1]), initials, fill=(255, 255, 255, 255), font=font)
    img.save(dest, "PNG")
    return dest


def _try_linkedin_photo(linkedin_url, dest, name):
    if not linkedin_url:
        return None
    try:
        req = urllib.request.Request(linkedin_url, headers={
            "User-Agent": "facebookexternalhit/1.1",
            "Accept": "text/html",
        })
        with urllib.request.urlopen(req, timeout=10) as resp:
            html = resp.read().decode("utf-8", errors="ignore")
        import re
        m = re.search(r'<meta[^>]+property=["\']og:image["\'][^>]+content=["\']([^"\'>]+)["\']', html)
        if not m:
            m = re.search(r'<meta[^>]+content=["\']([^"\'>]+)["\'][^>]+property=["\']og:image["\']', html)
        if m:
            img_url = m.group(1)
            if "media.licdn.com" in img_url or "static.licdn.com" in img_url:
                result = _download_image(img_url, dest, f"LinkedIn({name})")
                if result:
                    return result
    except Exception as e:
        print(f"  LinkedIn scrape failed for {name}: {e}")
    return None


def _fetch_speaker_photo(name, email=None, photo=None, linkedin=None, output_dir=None):
    import hashlib
    if output_dir is None:
        output_dir = tempfile.gettempdir()
    safe_name = name.replace(" ", "_").lower()
    dest = os.path.join(output_dir, f"speaker_{safe_name}.png")

    if photo:
        if os.path.exists(photo):
            from PIL import Image as _Img
            _Img.open(photo).verify()
            return photo
        result = _download_image(photo, dest, f"Speaker({name})")
        if result:
            return result

    result = _try_linkedin_photo(linkedin, dest, name)
    if result:
        return result

    if not email:
        parts = name.strip().lower().split()
        if len(parts) >= 2:
            email = f"{parts[0]}.{parts[-1]}@snowflake.com"

    if email:
        h = hashlib.sha256(email.strip().lower().encode()).hexdigest()
        check_url = f"https://gravatar.com/avatar/{h}?s=400&d=404"
        result = _download_image(check_url, dest, f"Gravatar({name})")
        if result:
            return result

    return _generate_initials_avatar(name, dest)


def _crop_circle(image_path):
    from PIL import Image as _Img, ImageDraw as _Draw
    img = _Img.open(image_path).convert("RGBA")
    size = min(img.size)
    left = (img.width - size) // 2
    top = (img.height - size) // 2
    img = img.crop((left, top, left + size, top + size))
    mask = _Img.new("L", (size, size), 0)
    draw = _Draw.Draw(mask)
    draw.ellipse((0, 0, size, size), fill=255)
    output = _Img.new("RGBA", (size, size), (0, 0, 0, 0))
    output.paste(img, mask=mask)
    circle_path = image_path.replace(".png", "_circle.png")
    output.save(circle_path, "PNG")
    return circle_path


SPEAKER_PHOTO_SIZE = Emu(1143000)
SPEAKER_NAME_FONT = Pt(14)
SPEAKER_TITLE_FONT = Pt(11)
SPEAKER_AREA_TOP = Emu(1600000)
SPEAKER_TEXT_GAP = Emu(114300)
SPEAKER_COL_WIDTH = Emu(2286000)
SPEAKER_COL_GAP = Emu(457200)


def apply_speakers(slide, data):
    shapes = find_shapes_with_text(slide)
    title_set = False
    subtitle_set = False
    shapes_to_remove = []
    shapes_by_top = sorted(shapes, key=lambda s: s.top)

    for shape in shapes_by_top:
        page_num = shape.text_frame.text.strip()
        if page_num == '\u2039#\u203a' or (shape.width < Emu(600000) and shape.top > Emu(4500000)):
            continue
        if shape.top < Emu(500000) and shape.width > Emu(4000000) and not title_set:
            set_shape_text(shape, data.get("title", "Your Snowflake Team"))
            constrain_title_autofit(shape)
            title_set = True
        elif Emu(500000) <= shape.top < Emu(1200000) and not subtitle_set:
            subtitle_text = data.get("subtitle", "")
            if subtitle_text:
                set_shape_text(shape, subtitle_text)
                constrain_title_autofit(shape)
            else:
                shapes_to_remove.append(shape)
            subtitle_set = True
        elif shape.top >= Emu(1200000) and shape.width > Emu(4000000):
            shapes_to_remove.append(shape)

    for s in shapes_to_remove:
        remove_shape(slide, s)

    speakers = data.get("speakers", [])
    if not speakers:
        set_speaker_notes(slide, data.get("speaker_notes", ""))
        return

    output_dir = tempfile.gettempdir()
    num = len(speakers)
    col_w = int(SPEAKER_COL_WIDTH)
    gap = int(SPEAKER_COL_GAP)
    total_block = num * col_w + (num - 1) * gap
    block_left = (int(SLIDE_WIDTH) - total_block) // 2

    for i, sp in enumerate(speakers):
        col_left = block_left + i * (col_w + gap)
        photo_left = col_left + (col_w - int(SPEAKER_PHOTO_SIZE)) // 2

        photo_path = _fetch_speaker_photo(
            sp.get("name", ""),
            email=sp.get("email"),
            photo=sp.get("photo"),
            linkedin=sp.get("linkedin"),
            output_dir=output_dir,
        )
        if photo_path:
            circle_path = _crop_circle(photo_path)
            slide.shapes.add_picture(
                circle_path,
                Emu(photo_left),
                SPEAKER_AREA_TOP,
                width=SPEAKER_PHOTO_SIZE,
                height=SPEAKER_PHOTO_SIZE,
            )

        name_top = int(SPEAKER_AREA_TOP) + int(SPEAKER_PHOTO_SIZE) + int(SPEAKER_TEXT_GAP)
        txBox = slide.shapes.add_textbox(
            Emu(col_left), Emu(name_top), Emu(col_w), Emu(400000)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = 1
        run = p.add_run()
        run.text = sp.get("name", "")
        run.font.size = SPEAKER_NAME_FONT
        run.font.bold = True
        run.font.color.rgb = MEDIUM_GRAY

        title_text = sp.get("title", "")
        if title_text:
            p2 = tf.add_paragraph()
            p2.alignment = 1
            run2 = p2.add_run()
            run2.text = title_text
            run2.font.size = SPEAKER_TITLE_FONT
            run2.font.color.rgb = MEDIUM_GRAY

    set_speaker_notes(slide, data.get("speaker_notes", ""))


APPLY_FN = {
    "cover": apply_cover,
    "cover_logo": apply_cover_logo,
    "safe_harbor": lambda slide, data: set_speaker_notes(slide, data.get("speaker_notes", "")),
    "agenda": apply_agenda,
    "chapter": apply_chapter,
    "content_1col": apply_content_1col,
    "content_2col_titled": apply_content_2col_titled,
    "content_3col_titled": apply_content_3col_titled,
    "content_3col_icons": apply_content_3col_icons,
    "content_4col_icons": apply_content_4col_icons,
    "template_clone": apply_template_clone,
    "speakers": apply_speakers,
    "thank_you": apply_thank_you,
}

ICON_APPLY_TYPES = {"content_3col_icons", "content_4col_icons"}


def template_index_for_type(slide_type, variant=None, template_index=None):
    if slide_type == "template_clone" and template_index is not None:
        return template_index
    if slide_type == "chapter" and variant:
        key = f"chapter_{variant}"
        return TEMPLATE_SLIDE_MAP.get(key, TEMPLATE_SLIDE_MAP["chapter_1"])
    if slide_type == "cover_logo":
        return TEMPLATE_SLIDE_MAP["cover_logo"]
    return TEMPLATE_SLIDE_MAP.get(slide_type, TEMPLATE_SLIDE_MAP["content_1col"])


def delete_slide(prs, slide_index):
    rId = prs.slides._sldIdLst[slide_index].get(qn("r:id"))
    prs.part.drop_rel(rId)
    sldId = prs.slides._sldIdLst[slide_index]
    prs.slides._sldIdLst.remove(sldId)


def _download_image(url, dest, label=""):
    try:
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=10) as resp:
            data = resp.read()
            if len(data) < 200:
                return None
            with open(dest, "wb") as f:
                f.write(data)
            from PIL import Image as _Img
            _Img.open(dest).verify()
            print(f"Logo fetched via {label} -> {dest}")
            return dest
    except Exception as e:
        print(f"Logo fetch failed ({label}): {e}", file=sys.stderr)
        if os.path.exists(dest):
            os.remove(dest)
        return None


MIN_LOGO_PX = 64
PREFERRED_LOGO_PX = 256


def _scrape_website_logo(domain, dest):
    import re
    try:
        url = f"https://www.{domain}"
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=10) as resp:
            html = resp.read().decode("utf-8", errors="ignore")
    except Exception as e:
        print(f"Website scrape failed for {domain}: {e}", file=sys.stderr)
        return None

    candidates = []

    og_match = re.search(r'(?:property|name)=["\']og:image["\'][^>]*content=["\']([^"\']+)["\']', html, re.I)
    if not og_match:
        og_match = re.search(r'content=["\']([^"\']+)["\'][^>]*(?:property|name)=["\']og:image["\']', html, re.I)
    if og_match:
        candidates.append(og_match.group(1))

    logo_urls = re.findall(r'(?:src|href)=["\']([^"\']*logo[^"\']*\.(?:png|svg|jpg|jpeg|webp))["\']', html, re.I)
    candidates.extend(logo_urls)

    img_logo_urls = re.findall(r'(?:src|href)=["\']([^"\']+\.(?:png|svg|jpg|jpeg|webp))["\']', html, re.I)
    for u in img_logo_urls:
        if "logo" in u.lower() and u not in candidates:
            candidates.append(u)

    for candidate_url in candidates:
        if candidate_url.startswith("//"):
            candidate_url = "https:" + candidate_url
        elif candidate_url.startswith("/"):
            candidate_url = f"https://www.{domain}{candidate_url}"
        elif not candidate_url.startswith("http"):
            continue

        if candidate_url.lower().endswith(".svg"):
            continue

        result = _download_image(candidate_url, dest, f"Website({domain})")
        if result:
            try:
                from PIL import Image as _Img
                img = _Img.open(result)
                w, h = img.size
                if w >= PREFERRED_LOGO_PX or h >= PREFERRED_LOGO_PX:
                    print(f"High-res logo found: {w}x{h}px from {candidate_url}")
                    return result
                else:
                    print(f"Logo too small from website ({w}x{h}px), trying next...", file=sys.stderr)
            except Exception:
                pass

    return None


def _check_logo_quality(path):
    try:
        from PIL import Image as _Img
        img = _Img.open(path)
        w, h = img.size
        if w < MIN_LOGO_PX or h < MIN_LOGO_PX:
            print(f"Logo too small ({w}x{h}px, need >={MIN_LOGO_PX}px): {path}", file=sys.stderr)
            return False
        return True
    except Exception as e:
        print(f"Logo quality check failed: {e}", file=sys.stderr)
        return False


def _convert_to_white(path):
    from PIL import Image as _Img
    import math
    img = _Img.open(path).convert("RGBA")
    w, h = img.size

    corners = [
        img.getpixel((0, 0)),
        img.getpixel((w - 1, 0)),
        img.getpixel((0, h - 1)),
        img.getpixel((w - 1, h - 1)),
    ]

    transparent_corners = sum(1 for c in corners if c[3] < 30)
    if transparent_corners >= 3:
        bg_r, bg_g, bg_b = 0, 0, 0
        bg_is_transparent = True
    else:
        opaque = [c for c in corners if c[3] >= 30]
        bg_r = sum(c[0] for c in opaque) // len(opaque)
        bg_g = sum(c[1] for c in opaque) // len(opaque)
        bg_b = sum(c[2] for c in opaque) // len(opaque)
        bg_is_transparent = False

    is_icon = abs(w - h) <= 10 and max(w, h) <= 300
    BG_TOLERANCE = 60
    data = img.getdata()

    if is_icon:
        if bg_is_transparent:
            fg_pixels = [(r, g, b) for r, g, b, a in data if a >= 30]
        else:
            fg_pixels = [
                (r, g, b) for r, g, b, a in data
                if a >= 30 and math.sqrt((r - bg_r)**2 + (g - bg_g)**2 + (b - bg_b)**2) >= BG_TOLERANCE
            ]

        if fg_pixels:
            bright_px = [p for p in fg_pixels if 0.299*p[0] + 0.587*p[1] + 0.114*p[2] > 230]
            non_bright_px = [p for p in fg_pixels if 0.299*p[0] + 0.587*p[1] + 0.114*p[2] <= 180]
            mid_px = [p for p in fg_pixels if 180 < 0.299*p[0] + 0.587*p[1] + 0.114*p[2] <= 230]
            bright_ratio = len(bright_px) / len(fg_pixels) if fg_pixels else 0
            non_bright_ratio = len(non_bright_px) / len(fg_pixels) if fg_pixels else 0
            mid_ratio = len(mid_px) / len(fg_pixels) if fg_pixels else 0

            bimodal = (
                bright_ratio > 0.05
                and non_bright_ratio > 0.3
                and mid_ratio < 0.15
                and len(bright_px) > 100
            )

            if bimodal:
                new_data = []
                for r, g, b, a in data:
                    if a < 30:
                        new_data.append((0, 0, 0, 0))
                        continue
                    if not bg_is_transparent:
                        dist = math.sqrt((r - bg_r)**2 + (g - bg_g)**2 + (b - bg_b)**2)
                        if dist < BG_TOLERANCE:
                            new_data.append((0, 0, 0, 0))
                            continue
                    lum = 0.299 * r + 0.587 * g + 0.114 * b
                    if lum > 180:
                        new_data.append((255, 255, 255, 255))
                    else:
                        new_data.append((0, 0, 0, 0))
                img.putdata(new_data)
                bbox = img.getbbox()
                if bbox:
                    img = img.crop(bbox)
                white_path = path.replace(".png", "_white.png")
                img.save(white_path, "PNG")
                print(f"Converted icon to white (extracted light elements) -> {white_path} ({img.width}x{img.height}px)")
                return white_path

    new_data = []
    for r, g, b, a in data:
        if a < 30:
            new_data.append((0, 0, 0, 0))
            continue
        dist = math.sqrt((r - bg_r)**2 + (g - bg_g)**2 + (b - bg_b)**2)
        if dist < BG_TOLERANCE:
            new_data.append((0, 0, 0, 0))
        else:
            new_data.append((255, 255, 255, 255))
    img.putdata(new_data)
    bbox = img.getbbox()
    if bbox:
        img = img.crop(bbox)
    white_path = path.replace(".png", "_white.png")
    img.save(white_path, "PNG")
    print(f"Converted to white logo -> {white_path} ({img.width}x{img.height}px)")
    return white_path


def fetch_logo(domain, output_dir=None):
    if output_dir is None:
        output_dir = tempfile.gettempdir()
    dest = os.path.join(output_dir, f"{domain.replace('.', '_')}_logo.png")

    result = _scrape_website_logo(domain, dest)
    if result and _check_logo_quality(result):
        return _convert_to_white(result)

    result = _download_image(
        f"https://logo.uplead.com/{domain}", dest, "Uplead"
    )
    if result and _check_logo_quality(result):
        return _convert_to_white(result)

    result = _download_image(
        f"https://logo.clearbit.com/{domain}?size=512", dest, "Clearbit"
    )
    if result and _check_logo_quality(result):
        return _convert_to_white(result)

    result = _download_image(
        f"https://t2.gstatic.com/faviconV2?client=SOCIAL&type=FAVICON&fallback_opts=TYPE,SIZE,URL&url=https://{domain}&size=256",
        dest,
        "Google Favicon",
    )
    if result and _check_logo_quality(result):
        return _convert_to_white(result)

    print(f"WARNING: Could not fetch logo for {domain} from any source", file=sys.stderr)
    return None


def build_deck(template_path, slides_json, output_path, logo_path=None, logo_domain=None):
    with open(slides_json, "r") as f:
        slides_data = json.load(f)

    if not logo_path and not logo_domain:
        if isinstance(slides_data, list) and slides_data:
            logo_domain = slides_data[0].get("logo_domain")

    if not logo_path and logo_domain:
        logo_path = fetch_logo(logo_domain, output_dir=os.path.dirname(output_path) or None)

    if not logo_path:
        print("Warning: No logo provided or fetched. Deck will not have a customer logo.", file=sys.stderr)

    output_prs = Presentation(template_path)
    original_count = len(output_prs.slides)
    icon_index = _build_icon_index(output_prs)

    first_slide = True
    logo_cache = {}
    output_dir = os.path.dirname(output_path) or None
    for slide_def in slides_data:
        slide_type = slide_def.get("type", "content_1col")
        variant = slide_def.get("variant")
        tmpl_idx = template_index_for_type(slide_type, variant, template_index=slide_def.get("template_index"))

        if tmpl_idx >= original_count:
            print(f"Warning: template index {tmpl_idx} out of range, skipping", file=sys.stderr)
            continue

        new_slide = clone_slide(output_prs, tmpl_idx)

        apply_fn = APPLY_FN.get(slide_type)
        if apply_fn:
            if slide_type in ICON_APPLY_TYPES:
                apply_fn(new_slide, slide_def, icon_index=icon_index)
            else:
                apply_fn(new_slide, slide_def)

        per_slide_logo_domain = slide_def.get("logo_domain")
        if per_slide_logo_domain:
            if per_slide_logo_domain not in logo_cache:
                logo_cache[per_slide_logo_domain] = fetch_logo(per_slide_logo_domain, output_dir=output_dir)
            per_logo = logo_cache[per_slide_logo_domain]
            if per_logo:
                place_logo_top_right(new_slide, per_logo)
        elif first_slide and logo_path:
            place_logo_top_right(new_slide, logo_path)
            first_slide = False

        icon_label = slide_def.get("icon")
        if icon_label and slide_type == "content_1col":
            icon_group = _find_icon(icon_index, icon_label)
            if icon_group:
                icon_left = SLIDE_WIDTH - ICON_MARGIN_RIGHT - ICON_SIZE
                icon_top = FOOTER_TOP - ICON_MARGIN_BOTTOM - ICON_SIZE
                place_icon(new_slide, icon_group, icon_left, icon_top)

    for i in range(original_count - 1, -1, -1):
        delete_slide(output_prs, i)

    output_prs.save(output_path)
    print(f"Deck saved to {output_path} ({len(slides_data)} slides)")


def main():
    parser = argparse.ArgumentParser(description="Generate branded Snowflake PPTX deck")
    parser.add_argument("--template", required=True, help="Path to Snowflake template PPTX")
    parser.add_argument("--slides-json", required=True, help="Path to slides content JSON")
    parser.add_argument("--output", required=True, help="Output PPTX path")
    parser.add_argument("--logo", default=None, help="Path to customer logo PNG")
    parser.add_argument("--logo-domain", default=None, help="Company domain for Clearbit logo fetch (e.g. brooksrunning.com)")
    args = parser.parse_args()

    build_deck(args.template, args.slides_json, args.output, logo_path=args.logo, logo_domain=args.logo_domain)


if __name__ == "__main__":
    main()
